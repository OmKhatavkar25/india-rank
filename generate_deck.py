"""Generate the hackathon submission deck as a PPTX file."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x96, 0xD6)
GREEN = RGBColor(0x00, 0xC8, 0x88)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from pptx.oxml.ns import qn
        solidFill = shape.fill._fill
        srgbClr = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        if srgbClr is not None:
            alpha_el = srgbClr.find(qn('a:alpha'))
            if alpha_el is None:
                from lxml import etree
                alpha_el = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha_el.set('val', str(int(alpha * 1000)))
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_slide(slide, left, top, width, height, items, font_size=16, color=BLACK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "Redrob India Runs — Data & AI Challenge", 36, True, WHITE, PP_ALIGN.LEFT)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
             "Intelligent Candidate Discovery & Ranking System", 24, False, RGBColor(0xAA, 0xCC, 0xEE), PP_ALIGN.LEFT)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
             "Ranking 100K candidates for Senior AI Engineer — Founding Team @ Redrob AI", 18, False, RGBColor(0x88, 0xAA, 0xCC), PP_ALIGN.LEFT)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
             "Team: Om Khatavkar", 16, False, RGBColor(0x88, 0xAA, 0xCC), PP_ALIGN.LEFT)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "GitHub: github.com/omkhatavkar/candidate-ranker", 14, False, GRAY, PP_ALIGN.LEFT)


# ============================================================
# SLIDE 2: PROBLEM STATEMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.04), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "The Problem", 32, True, BLACK)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
             "Recruiters miss the right person — keyword filters can't see what actually matters.", 20, False, GRAY)

items = [
    "▸  A recruiter reviews hundreds of profiles for a single role",
    "▸  Keyword-based filters miss candidates with adjacent/transferable skills",
    "▸  Static profiles don't capture availability, responsiveness, or intent",
    "▸  Consulting-firm candidates flood the pool but don't fit product roles",
    "▸  Honeypot profiles with impossible experience combos waste recruiter time",
]
add_bullet_slide(slide, Inches(0.8), Inches(2.4), Inches(11), Inches(4.0), items, 18, BLACK)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
             "Goal: Rank candidates like a great recruiter would — understand context, not just keywords", 16, True, ACCENT)


# ============================================================
# SLIDE 3: THE JD — SENIOR AI ENGINEER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.04), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Understanding the Job Description", 32, True, BLACK)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
             "Senior AI Engineer — Founding Team @ Redrob AI (Series A)", 20, False, GRAY)

# Left column — what the JD says
add_shape(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(4.5), RGBColor(0xE8, 0xF4, 0xFD))
add_text_box(slide, Inches(1.0), Inches(2.6), Inches(5.0), Inches(0.5),
             "What the JD explicitly asks for", 18, True, ACCENT)
jd_items = [
    "• Production ML experience (embeddings, retrieval, ranking)",
    "• Vector DB/hybrid search infrastructure",
    "• Strong Python and code quality",
    "• Evaluation frameworks (NDCG, MRR, MAP)",
    "• Product company experience (not consulting)",
    "• 5-9 years experience sweet spot",
    "• Pune/Noida location preferred",
    "• < 30 day notice period preferred",
]
add_bullet_slide(slide, Inches(1.0), Inches(3.3), Inches(5.0), Inches(3.5), jd_items, 15, BLACK, Pt(4))

# Right column — traps
add_shape(slide, Inches(7.0), Inches(2.4), Inches(5.5), Inches(4.5), RGBColor(0xFD, 0xE8, 0xE8))
add_text_box(slide, Inches(7.2), Inches(2.6), Inches(5.0), Inches(0.5),
             "What the JD warns against", 18, True, RGBColor(0xCC, 0x33, 0x33))
trap_items = [
    "✗ Don't count AI keywords — that's a trap",
    "✗ Marketing Manager with ML skills listed ≠ fit",
    "✗ Consulting-only backgrounds (TCS, Infosys, Wipro...)",
    "✗ Title-chasers (avg tenure < 18 months)",
    "✗ LangChain tutorial profiles",
    "✗ Pure research without production deployment",
    "✗ ~80 honeypot profiles with impossible data",
]
add_bullet_slide(slide, Inches(7.2), Inches(3.3), Inches(5.0), Inches(3.5), trap_items, 15, BLACK, Pt(4))


# ============================================================
# SLIDE 4: ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.04), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "System Architecture", 32, True, BLACK)

# Pipeline flow boxes
boxes = [
    ("candidates.jsonl\n100K lines", 0.8, 2.0, 2.2, 1.2, RGBColor(0xE8, 0xF4, 0xFD), BLACK),
    ("Loader\n(JSONL parser)", 3.5, 2.0, 2.0, 1.2, RGBColor(0xD4, 0xED, 0xDA), BLACK),
    ("Feature Extraction\n(30+ features)", 6.0, 2.0, 2.2, 1.2, RGBColor(0xFD, 0xF0, 0xE0), BLACK),
    ("Scoring Engine\n(7 dimensions)", 8.7, 2.0, 2.2, 1.2, RGBColor(0xE8, 0xF4, 0xFD), BLACK),
    ("Ranker\n(Top 100)", 11.4, 2.0, 1.5, 1.2, RGBColor(0x00, 0x96, 0xD6), WHITE),
]

for text, x, y, w, h, bg, fc in boxes:
    add_shape(slide, Inches(x), Inches(y), Inches(w), Inches(h), bg)
    add_text_box(slide, Inches(x + 0.1), Inches(y + 0.1), Inches(w - 0.2), Inches(h - 0.2),
                 text, 14, True, fc, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x + 0.1), Inches(y + h + 0.05), Inches(w - 0.2), Inches(0.3),
                 "→", 20, True, ACCENT, PP_ALIGN.CENTER)

# Detailed feature breakdown
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.5),
             "Feature Extraction — 30+ signals per candidate", 20, True, BLACK)

feat_items = [
    "Role History:  ML/AI title detection, retrieval/ranking keyword search in descriptions, production deployment evidence",
    "Company Analysis:  Product company vs consulting-firm classification (TCS/Infosys/Wipro/Accenture blacklist)",
    "Skills:  JD-core skills match (embeddings, vector DBs, Python, fine-tuning, NDCG), weighted by proficiency × endorsements × duration",
    "Behavioral:  Recruiter response rate, last active date, profile views, saves by recruiters, GitHub activity, notice period",
    "Stability:  Average tenure, job-hopping detection (<12 months at any role), career progression",
    "Education & Location:  Tier score, field-of-study relevance, city match (Pune/Noida/Bangalore/Hyderabad/Delhi)",
    "Honeypot Detection:  Expert proficiency with <6mo duration, 20+ skills averaging <6mo, suspicious patterns",
]
add_bullet_slide(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(3.0), feat_items, 15, BLACK, Pt(3))


# ============================================================
# SLIDE 5: SCORING BREAKDOWN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.04), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Scoring Dimensions & Weights", 32, True, BLACK)

# Weight bars
dims = [
    ("ML/AI Experience", "30%", 0.30, RGBColor(0x00, 0x96, 0xD6)),
    ("Product Company", "20%", 0.20, RGBColor(0x00, 0x96, 0xD6)),
    ("Skill Relevance", "20%", 0.20, RGBColor(0x00, 0x96, 0xD6)),
    ("Career Stability", "10%", 0.10, RGBColor(0x00, 0x96, 0xD6)),
    ("Behavioral Signals", "10%", 0.10, RGBColor(0x00, 0x96, 0xD6)),
    ("Experience Years", "5%", 0.05, RGBColor(0x00, 0x96, 0xD6)),
    ("Education & Location", "5%", 0.05, RGBColor(0x00, 0x96, 0xD6)),
]

y_start = 1.8
bar_max_width = 8.0
for i, (label, pct, weight, color) in enumerate(dims):
    y = y_start + i * 0.7
    add_text_box(slide, Inches(0.8), Inches(y), Inches(2.5), Inches(0.5),
                 label, 15, False, BLACK)
    add_text_box(slide, Inches(3.4), Inches(y), Inches(0.8), Inches(0.5),
                 pct, 14, True, ACCENT, PP_ALIGN.RIGHT)
    bar_width = bar_max_width * weight
    add_shape(slide, Inches(4.3), Inches(y + 0.05), Inches(bar_width), Inches(0.35), color)
    if weight > 0.1:
        add_text_box(slide, Inches(4.3), Inches(y + 0.05), Inches(bar_width), Inches(0.35),
                     f"  {label.lower()}: what it means", 11, False, WHITE, PP_ALIGN.LEFT)

add_text_box(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
             "Total = weighted sum of all dimensions, clamped to [0, 1], minus honeypot penalties", 14, False, GRAY)


# ============================================================
# SLIDE 6: KEY INSIGHTS — WHY NOT KEYWORD MATCHING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.04), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Why Not Simple Keyword Matching?", 32, True, BLACK)

# Three columns
col_data = [
    ("Trap: Keyword Counters", [
        "A Marketing Manager with",
        "'Machine Learning' in skills",
        "",
        "Score if keyword counting:",
        "  HIGH (has AI keywords)",
        "",
        "Score with our system:",
        "  LOW (title mismatch,",
        "  no ML career history,",
        "  non-product industry)",
    ], RGBColor(0xFD, 0xE8, 0xE8), RGBColor(0xCC, 0x33, 0x33)),
    ("Missed: Tier 5 Candidates", [
        "A Data Scientist whose",
        "profile never says 'RAG'",
        "or 'Pinecone' but who",
        "built a recommendation",
        "system at a product company.",
        "",
        "Score if keyword counting:",
        "  MEDIUM (no buzzwords)",
        "",
        "Score with our system:",
        "  HIGH (actual relevant",
        "  career history detected)",
    ], RGBColor(0xE8, 0xF4, 0xFD), ACCENT),
    ("Honeypots", [
        "Profiles with impossible",
        "combinations:",
        "• 'Expert' in 10 skills",
        "  with 0 years experience",
        "• Experience at a company",
        "  founded after they",
        "  supposedly worked there",
        "",
        "Our system detects:",
        "  suspicious skill patterns",
        "  and penalizes them",
    ], RGBColor(0xE8, 0xF4, 0xFD), RGBColor(0xCC, 0x66, 0x00)),
]

for i, (title, items, bg, accent) in enumerate(col_data):
    x = 0.8 + i * 4.2
    add_shape(slide, Inches(x), Inches(1.8), Inches(3.8), Inches(5.2), bg)
    add_text_box(slide, Inches(x + 0.3), Inches(2.0), Inches(3.2), Inches(0.5),
                 title, 16, True, accent)
    add_bullet_slide(slide, Inches(x + 0.3), Inches(2.6), Inches(3.2), Inches(4.2), items, 14, BLACK, Pt(1))


# ============================================================
# SLIDE 7: RESULTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.04), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Results & Validation", 32, True, BLACK)

result_items = [
    "▸ Processed 100,000 candidates in 11.5 seconds on CPU (8 cores, 16 GB RAM)",
    "▸ No GPU required, no network calls, well within the 5-minute limit",
    "▸ Output: 100 rows, validated by validate_submission.py — all checks pass",
    "",
    "Top ranked candidates share these traits:",
    "  • 4-8 years of actual ML/AI engineering experience at product companies",
    "  • Retrieval, ranking, or search system experience in their career history",
    "  • 2-10 JD-relevant skills (embeddings, vector DBs, Python, fine-tuning, etc.)",
    "  • High recruiter response rate (>70%), recently active on platform",
    "  • Located in Pune, Noida, Bangalore, Hyderabad, Delhi, or Gurgaon",
    "  • Short notice period (<30 days preferred)",
    "",
    "Sample output (Rank #1):",
    "  CAND_0081846 | score: 0.8799 | reasoning: \"6.6yrs ML/AI experience;",
    "  retrieval/ranking/systems exp; product company background; 5 relevant skills;",
    "  6.7yrs exp (ideal range); response rate 73%; short notice\"",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5), result_items, 16, BLACK, Pt(3))


# ============================================================
# SLIDE 8: SUBMISSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)
add_shape(slide, Inches(0), Inches(1.2), Inches(13.333), Inches(0.04), ACCENT)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Submission Checklist", 32, True, BLACK)

submit_items = [
    "✓ submission.csv — 100 ranked candidates in required format (validated)",
    "✓ GitHub repository — clean, complete, working code at github.com/omkhatavkar/candidate-ranker",
    "✓ submission_metadata.yaml — team info, compute env, AI tools declaration",
    "✓ This deck — explaining approach, architecture, and design decisions",
    "",
    "Reproduce command:",
    "  python rank.py --candidates ./candidates.jsonl --out ./submission.csv",
    "",
    "Compute environment:",
    "  • Python 3.11+ | 8 CPU cores | 16 GB RAM | Windows",
    "  • No GPU | No network during ranking | ~12 seconds runtime",
    "",
]
add_bullet_slide(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.5), submit_items, 16, BLACK, Pt(3))


# ============================================================
# SLIDE 9: THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape(slide, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1),
             "Thank You", 44, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(1),
             "Om Khatavkar  |  github.com/omkhatavkar/candidate-ranker", 20, False, RGBColor(0x88, 0xAA, 0xCC), PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
             "Making hiring smarter — one rank at a time.", 16, False, GRAY, PP_ALIGN.CENTER)


# Save
prs.save("D:\\india runs\\candidate-ranker\\submission_deck.pptx")
print("Deck saved to submission_deck.pptx")
